"""
Document processing utilities for extracting text from various file formats.
"""
import os
import json
import csv
import html
import re
from typing import List, Optional
from PyPDF2 import PdfReader
from pathlib import Path
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class DocumentProcessor:
    """Handles document processing and text extraction."""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from a PDF file."""
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading PDF: {str(e)}")
    
    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """Extract text from a text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            raise Exception(f"Error reading text file: {str(e)}")
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from a DOCX file."""
        if not DOCX_AVAILABLE:
            raise Exception("python-docx library is not installed. Install it with: pip install python-docx")
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading DOCX file: {str(e)}")
    
    @staticmethod
    def extract_text_from_markdown(file_path: str) -> str:
        """Extract text from a Markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            raise Exception(f"Error reading Markdown file: {str(e)}")
    
    @staticmethod
    def extract_text_from_csv(file_path: str) -> str:
        """Extract text from a CSV file."""
        try:
            text = ""
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row in reader:
                    text += " | ".join(row) + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading CSV file: {str(e)}")
    
    @staticmethod
    def extract_text_from_json(file_path: str) -> str:
        """Extract text from a JSON file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                # Convert JSON to readable text format
                return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception as e:
            raise Exception(f"Error reading JSON file: {str(e)}")
    
    @staticmethod
    def extract_text_from_html(file_path: str) -> str:
        """Extract text from an HTML file."""
        try:
            try:
                from bs4 import BeautifulSoup
                with open(file_path, 'r', encoding='utf-8') as f:
                    soup = BeautifulSoup(f.read(), 'html.parser')
                    # Remove script and style elements
                    for script in soup(["script", "style"]):
                        script.decompose()
                    # Get text and clean it up
                    text = soup.get_text()
                    # Clean up whitespace
                    lines = (line.strip() for line in text.splitlines())
                    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                    text = '\n'.join(chunk for chunk in chunks if chunk)
                    return text
            except ImportError:
                # Fallback without BeautifulSoup - basic HTML tag removal
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    # Remove HTML tags using regex (basic approach)
                    text = re.sub(r'<[^>]+>', '', content)
                    # Decode HTML entities
                    text = html.unescape(text)
                    return text.strip()
        except Exception as e:
            raise Exception(f"Error reading HTML file: {str(e)}")
    
    @staticmethod
    def extract_text_from_rtf(file_path: str) -> str:
        """Extract text from an RTF file."""
        try:
            import striprtf
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
                text = striprtf.striprtf.StripRtf().strip(rtf_content)
                return text.strip()
        except ImportError:
            # Fallback: basic RTF text extraction (removes RTF control codes)
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                # Remove RTF control codes (basic approach)
                text = re.sub(r'\\[a-z]+\d*\s?', '', content)
                text = re.sub(r'\{[^}]*\}', '', text)
                return text.strip()
        except Exception as e:
            raise Exception(f"Error reading RTF file: {str(e)}")
    
    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> str:
        """Extract text from an Excel file (.xlsx)."""
        if not EXCEL_AVAILABLE:
            raise Exception("openpyxl library is not installed. Install it with: pip install openpyxl")
        try:
            workbook = openpyxl.load_workbook(file_path)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n--- Sheet: {sheet_name} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text += row_text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading Excel file: {str(e)}")
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from a PowerPoint file (.pptx)."""
        if not PPTX_AVAILABLE:
            raise Exception("python-pptx library is not installed. Install it with: pip install python-pptx")
        try:
            prs = Presentation(file_path)
            text = ""
            for i, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slide {i} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading PowerPoint file: {str(e)}")
    
    @staticmethod
    def process_file(file_path: str) -> str:
        """Process a file and extract text based on file extension."""
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext == '.pdf':
            return DocumentProcessor.extract_text_from_pdf(file_path)
        elif file_ext == '.txt':
            return DocumentProcessor.extract_text_from_txt(file_path)
        elif file_ext == '.docx':
            return DocumentProcessor.extract_text_from_docx(file_path)
        elif file_ext in ['.md', '.markdown']:
            return DocumentProcessor.extract_text_from_markdown(file_path)
        elif file_ext == '.csv':
            return DocumentProcessor.extract_text_from_csv(file_path)
        elif file_ext == '.json':
            return DocumentProcessor.extract_text_from_json(file_path)
        elif file_ext in ['.html', '.htm']:
            return DocumentProcessor.extract_text_from_html(file_path)
        elif file_ext == '.rtf':
            return DocumentProcessor.extract_text_from_rtf(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            return DocumentProcessor.extract_text_from_xlsx(file_path)
        elif file_ext == '.pptx':
            return DocumentProcessor.extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}. Supported formats: .pdf, .txt, .docx, .md, .csv, .json, .html, .rtf, .xlsx, .pptx")
    
    @staticmethod
    def split_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
        """Split text into chunks for processing."""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            
            # Try to break at sentence boundary
            if end < len(text):
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                break_point = max(last_period, last_newline)
                
                if break_point > chunk_size * 0.5:  # Only break if we're past halfway
                    chunk = chunk[:break_point + 1]
                    end = start + break_point + 1
            
            chunks.append(chunk.strip())
            start = end - chunk_overlap
        
        return chunks

